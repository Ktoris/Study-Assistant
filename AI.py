import streamlit as st
import json
import requests
from PyPDF2 import PdfReader
from pptx import Presentation

# ================== CONFIG ==================
API_KEY = st.secrets["OPENROUTER_API_KEY"]
st.write("API Key loaded?", bool(API_KEY))  # True if loaded correctly

MODEL = "deepseek/deepseek-chat-v3.1:free"
BASE_URL = "https://openrouter.ai/api/v1/chat/completions"

# ================== HELPER FUNCTION ==================
def ask_openrouter(prompt, notes):
    payload = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": prompt},
            {"role": "user", "content": "Notes:\n" + notes}
        ]
    }
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    response = requests.post(BASE_URL, headers=headers, data=json.dumps(payload))
    response.raise_for_status()
    data = response.json()
    return data["choices"][0]["message"]["content"]

# ================== PROMPTS ==================
QUIZ_PROMPT = """You are an expert teacher creating practice tests. 
I will provide you with a set of notes on a topic. 
Based on these notes, generate ONLY multiple-choice questions in valid JSON format.
Do not include explanations or extra text.

The JSON must look like this:

{
  "multiple_choice": [
    {
      "question": "string",
      "options": [
        "Option A text",
        "Option B text",
        "Option C text",
        "Option D text"
      ],
      "answer": "Option X text"
    }
  ]
}

Rules:
- Every question must have exactly 4 distinct options.
- The answer must exactly match one of the 4 options. Do not include ```json fences or any explanation.
"""

FEYNMAN_PROMPT = """You are an expert teacher using the Feynman technique. 
I will provide you with a set of notes. 
Explain the concepts in the simplest possible way, as if teaching a 12-year-old.
Use analogies, simple words, and short sentences.
Do not output JSON, just plain text. Do not include ```json fences or any explanation.
"""

PRACTICE_TEST_PROMPT = """You are an expert teacher creating practice tests. 
I will provide you with a set of notes. 
Return ONLY valid JSON. No markdown, no extra text.

The JSON must look like this:

{
  "practice_test": [
    {
      "type": "multiple_choice",
      "question": "string",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "answer": "Option X"
    },
    {
      "type": "true_false",
      "question": "string",
      "answer": true
    },
    {
      "type": "fill_blank",
      "question": "The capital of France is ____.",
      "answer": "Paris"
    },
    {
      "type": "open_question",
      "question": "Explain the causes of World War II."
    }
  ]
}

Rules:
- Multiple-choice must always have 4 options.
- True/False answers must be strictly true or false (boolean).
- Fill-in-the-blank must use '____' for the blank space.
- Open questions have no answer field.
- Do not include ```json fences or any explanation.
"""

SUMMARY_PROMPT = """You are an expert teacher. 
I will provide you with a set of notes. 
Generate a clear and concise summary of the notes in plain text. 
Rules:
- Keep the summary short (1–3 paragraphs).
- Use simple, clear sentences.
- Do not include lists unless necessary.
"""

# ================== PAGE SETUP ==================
st.set_page_config(page_title="AI Study App", page_icon="📘", layout="centered")

st.markdown("""
<style>
.main-title { text-align: center; font-size: 2.5em; margin-bottom: 0; }
.subtitle { text-align: center; color: gray; margin-top: 0; margin-bottom: 2em; }
.stTabs [role="tab"] { padding: 12px 20px; font-size: 1.05em; }
</style>
""", unsafe_allow_html=True)

st.markdown("<h1 class='main-title'>📘 AI-Powered Study App</h1>", unsafe_allow_html=True)
st.markdown("<p class='subtitle'>Upload notes, generate quizzes, simplify with Feynman, or get summaries — all powered by AI.</p>", unsafe_allow_html=True)

# ================== INPUT SECTION ==================
st.header("✍️ Input Notes")
notes = st.text_area("Paste your notes here:", height=200)

col1, col2 = st.columns(2)
with col1:
    uploaded_pdf = st.file_uploader("📂 Upload a PDF", type=["pdf"])
with col2:
    uploaded_ppt = st.file_uploader("📂 Upload a PowerPoint", type=["pptx"])

if uploaded_pdf:
    pdf_reader = PdfReader(uploaded_pdf)
    notes = "\n".join([page.extract_text() or "" for page in pdf_reader.pages]).strip()
    st.success("✅ PDF uploaded and converted to text!")

def extract_text_from_pptx(file):
    prs = Presentation(file)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]).strip()

if uploaded_ppt:
    notes = extract_text_from_pptx(uploaded_ppt)
    st.success("✅ PowerPoint uploaded and converted to text!")

# ================== FUNCTION CHOICES ==================
tabs = st.tabs(["📝 Quiz", "🧠 Feynman", "📑 Practice Test", "📖 Summary"])

# -------------------- QUIZ TAB --------------------
with tabs[0]:
    st.subheader("Generate a Quiz from Notes")
    if st.button("⚡ Create Quiz"):
        with st.spinner("Generating quiz..."):
            try:
                response_text = ask_openrouter(QUIZ_PROMPT, notes)
                quiz_data = json.loads(response_text)
                st.session_state.quiz = quiz_data["multiple_choice"]
                st.session_state.user_answers = {}
                st.success("Quiz generated successfully!")
            except Exception as e:
                st.error(f"Failed to generate quiz: {e}")

    if "quiz" in st.session_state:
        for i, q in enumerate(st.session_state.quiz):
            st.subheader(f"Q{i+1}: {q['question']}")
            choice = st.radio("Choose an answer:", q["options"], key=f"q{i}")
            st.session_state.user_answers[i] = choice

        if st.button("Submit Quiz"):
            correct = 0
            total = len(st.session_state.quiz)
            for i, q in enumerate(st.session_state.quiz):
                user_ans = st.session_state.user_answers.get(i)
                correct_ans = q["answer"]
                if user_ans == correct_ans:
                    st.success(f"Q{i+1}: ✅ Correct")
                    correct += 1
                else:
                    st.error(f"Q{i+1}: ❌ Wrong. Correct answer: {correct_ans}")
            st.subheader(f"Final Score: {correct} / {total}")

# -------------------- FEYNMAN TAB --------------------
with tabs[1]:
    st.subheader("Explain with Feynman Technique")
    if st.button("💡 Simplify Notes"):
        with st.spinner("Explaining..."):
            try:
                st.session_state.feynman_explanation = ask_openrouter(FEYNMAN_PROMPT, notes)
                st.success("Explanation generated!")
            except Exception as e:
                st.error(f"Failed: {e}")
    if "feynman_explanation" in st.session_state:
        st.write(st.session_state.feynman_explanation)

# -------------------- PRACTICE TEST TAB --------------------
with tabs[2]:
    st.subheader("Generate a Practice Test")
    if st.button("📝 Create Practice Test"):
        with st.spinner("Generating practice test..."):
            try:
                response_text = ask_openrouter(PRACTICE_TEST_PROMPT, notes)
                practice_data = json.loads(response_text)
                st.session_state.practice_test = practice_data["practice_test"]
                st.session_state.practice_answers = {}
                st.success("Practice test generated successfully!")
            except Exception as e:
                st.error(f"Failed: {e}")

    if "practice_test" in st.session_state:
        for i, q in enumerate(st.session_state.practice_test):
            q_type = q["type"]
            if q_type == "multiple_choice":
                st.session_state.practice_answers[i] = st.radio(q["question"], q["options"], key=f"pt_mc_{i}")
            elif q_type == "true_false":
                st.session_state.practice_answers[i] = st.radio(q["question"], [True, False], key=f"pt_tf_{i}")
            elif q_type == "fill_blank":
                st.session_state.practice_answers[i] = st.text_input(q["question"], key=f"pt_fb_{i}")
            elif q_type == "open_question":
                st.text_area(q["question"], key=f"pt_oq_{i}")

        if st.button("Submit Practice Test"):
            correct, total = 0, 0
            for i, q in enumerate(st.session_state.practice_test):
                q_type = q["type"]
                if q_type in ["multiple_choice", "true_false", "fill_blank"]:
                    total += 1
                    user_ans = st.session_state.practice_answers.get(i)
                    correct_ans = q["answer"]
                    if isinstance(correct_ans, str):
                        is_correct = str(user_ans).strip().lower() == correct_ans.strip().lower()
                    else:
                        is_correct = user_ans == correct_ans
                    if is_correct:
                        st.success(f"Q{i+1}: ✅ Correct")
                        correct += 1
                    else:
                        st.error(f"Q{i+1}: ❌ Wrong. Correct answer: {correct_ans}")
                elif q_type == "open_question":
                    st.info(f"Q{i+1}: Open question — check your own answer.")
            st.subheader(f"Final Score: {correct} / {total}")

# -------------------- SUMMARY TAB --------------------
with tabs[3]:
    st.subheader("Summarize Notes")
    if st.button("📖 Summarize"):
        with st.spinner("Summarizing..."):
            try:
                st.session_state.summary = ask_openrouter(SUMMARY_PROMPT, notes)
                st.success("Summary generated!")
            except Exception as e:
                st.error(f"Failed: {e}")
    if "summary" in st.session_state:
        st.write(st.session_state.summary)
